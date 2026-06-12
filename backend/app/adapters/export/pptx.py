"""Presentation model → .pptx (`python-pptx`).

One content slide per stored `slides` entry: headline as the slide title,
evidence bullets with `[n]` citation refs (numbered by first appearance across
the whole deck, shared with the references slide), a rendered `VisualSpec`
where one is stored, and speaker notes in the notes pane. Tabular visual types
(comparison_table / timeline / trend) render as native tables — deterministic
and robust to re-export; `bullet_set` renders as bullets.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation as PptxDocument
from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from app.adapters.export.base import (
    Exporter,
    PresentationContent,
    format_reference,
    presentation_source_ids,
)

_TITLE_LAYOUT = 0
_TITLE_AND_CONTENT_LAYOUT = 1

_VISUAL_DEFAULT_COLUMNS = {"timeline": ["When", "What"], "trend": ["x", "y"]}


def _refs(source_ids: Any, numbering: dict[str, int]) -> str:
    markers = "".join(
        f"[{numbering[s]}]" for s in (source_ids or []) if isinstance(s, str) and s in numbering
    )
    return f" {markers}" if markers else ""


def _set_bullets(slide: PptxSlide, lines: list[str]) -> None:
    body = slide.placeholders[1].text_frame
    body.word_wrap = True
    for i, line in enumerate(lines):
        paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(18)


def _add_visual_table(slide: PptxSlide, visual: dict[str, Any], top_inches: float) -> None:
    columns = [str(c) for c in visual.get("columns") or []]
    rows = [[str(cell) for cell in row] for row in visual.get("rows") or []]
    if not columns:
        width = max((len(r) for r in rows), default=2)
        columns = (_VISUAL_DEFAULT_COLUMNS.get(str(visual.get("type"))) or [])[:width]
        columns += [f"col{i + 1}" for i in range(len(columns), width)]
    if not columns:
        return
    n_rows, n_cols = len(rows) + 1, len(columns)
    height = Inches(min(0.4 * n_rows, 7.0 - top_inches))
    frame = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(top_inches), Inches(9.0), height
    )
    table = frame.table
    for c, header in enumerate(columns):
        table.cell(0, c).text = header
    for r, row in enumerate(rows, start=1):
        padded = (row + [""] * n_cols)[:n_cols]
        for c, value in enumerate(padded):
            table.cell(r, c).text = value


class PptxPresentationExporter(Exporter[PresentationContent]):
    extension = "pptx"
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def render(self, content: PresentationContent) -> bytes:
        deck = content.presentation
        numbering = {s: i + 1 for i, s in enumerate(presentation_source_ids(deck))}
        notes_by_slide = {note.get("slide"): note.get("notes") for note in deck.speaker_notes or []}

        pptx = PptxDocument()

        title_slide = pptx.slides.add_slide(pptx.slide_layouts[_TITLE_LAYOUT])
        title_slide.shapes.title.text = deck.through_line or "Research presentation"
        subtitle = title_slide.placeholders[1]
        subtitle.text = (
            f"{len(deck.key_messages or [])} key messages · {len(deck.slides or [])} slides"
        )

        if deck.key_messages:
            messages_slide = pptx.slides.add_slide(pptx.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
            messages_slide.shapes.title.text = "Key messages"
            _set_bullets(
                messages_slide,
                [
                    f"{m.get('message', '')}{_refs(m.get('source_ids'), numbering)}"
                    for m in deck.key_messages
                ],
            )

        for i, stored in enumerate(deck.slides or []):
            slide = pptx.slides.add_slide(pptx.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
            slide.shapes.title.text = stored.get("headline", "")
            bullets = []
            for point in stored.get("evidence") or []:
                inference = " (inference)" if point.get("is_inference") else ""
                bullets.append(
                    f"{point.get('text', '')}"
                    f"{_refs(point.get('source_ids'), numbering)}{inference}"
                )
            visual = stored.get("visual")
            if visual and visual.get("type") == "bullet_set":
                if visual.get("title"):
                    bullets.append(str(visual["title"]))
                bullets.extend(str(p) for p in visual.get("points") or [])
                visual = None
            if bullets:
                _set_bullets(slide, bullets)
            if visual:
                # Below the bullets when both are present; higher when alone.
                _add_visual_table(slide, visual, top_inches=4.2 if bullets else 2.0)
            notes = notes_by_slide.get(i)
            if notes:
                slide.notes_slide.notes_text_frame.text = str(notes)

        if numbering:
            refs_slide = pptx.slides.add_slide(pptx.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
            refs_slide.shapes.title.text = "References"
            _set_bullets(
                refs_slide,
                [
                    format_reference(number, source_id, content.sources.get(source_id))
                    for source_id, number in sorted(numbering.items(), key=lambda kv: kv[1])
                ],
            )

        buffer = io.BytesIO()
        pptx.save(buffer)
        return buffer.getvalue()
