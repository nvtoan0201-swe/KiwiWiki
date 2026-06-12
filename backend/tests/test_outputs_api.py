"""Report/presentation REST endpoints (phase 5): read, user edits as new
versions, the regenerating rewrite (acceptance 6), and file exports
(acceptance 5) — all through the HTTP surface."""

import io

import pytest
import pytest_asyncio
from docx import Document
from pptx import Presentation as PptxDocument
from sqlalchemy import select

from app.api.reports import get_llm_json
from app.core.constants import AuditActionType
from app.db.models import AuditLogEntry, Presentation, Provenance, Report
from app.events.bus import set_event_bus
from app.schemas.report import ReportOutline, ReportSection, SelfCheckResult
from tests.llm_fakes import FakeLLM
from tests.orchestrator_utils import make_project
from tests.report_utils import (
    QUESTION,
    clean_self_check,
    outline_responder,
    section_responder,
    seed_corpus,
)
from tests.test_runner import RecordingBus


@pytest.fixture
def event_bus() -> RecordingBus:
    """Override the conftest bus with a recording one so the tests can assert
    on `output_ready` emitted by the rewrite endpoint."""
    bus = RecordingBus()
    set_event_bus(bus)
    return bus


@pytest_asyncio.fixture
async def project_with_outputs(sessionmaker):
    """A committed project with the grounded corpus, one report citing a real
    source, and one presentation."""
    project = await make_project(sessionmaker, research_question=QUESTION, audience="expert")
    async with sessionmaker() as session:
        first, second = await seed_corpus(session, project.id)
        report = Report(
            project_id=project.id,
            audience="expert",
            content_markdown=(
                "# Findings\n\n"
                f"Transformers lead on short horizons.[^src:{first.id}] "
                "*(confidence: well established)*\n"
            ),
            self_check_result={"clean": True, "findings": []},
            stopping_criterion="saturation",
            version=1,
        )
        presentation = Presentation(
            project_id=project.id,
            through_line="Horizon decides the winner.",
            key_messages=[
                {"message": "M1", "source_ids": [first.id]},
                {"message": "M2", "source_ids": [second.id]},
                {"message": "M3", "source_ids": []},
            ],
            slides=[
                {
                    "headline": "H1",
                    "key_message_index": 0,
                    "evidence": [
                        {
                            "text": "E1",
                            "source_ids": [first.id],
                            "passage": "p",
                            "is_inference": False,
                        }
                    ],
                    "visual": None,
                }
            ],
            speaker_notes=[{"slide": 0, "notes": "n0"}],
            version=1,
        )
        session.add_all([report, presentation])
        await session.commit()
        return {
            "project": project,
            "report_id": report.id,
            "presentation_id": presentation.id,
            "source_ids": [first.id, second.id],
        }


async def test_get_and_list_outputs(client, project_with_outputs):
    project_id = project_with_outputs["project"].id
    listed = await client.get(f"/projects/{project_id}/reports")
    assert listed.status_code == 200
    assert [r["version"] for r in listed.json()] == [1]

    got = await client.get(f"/reports/{project_with_outputs['report_id']}")
    assert got.status_code == 200
    assert got.json()["audience"] == "expert"

    decks = await client.get(f"/projects/{project_id}/presentations")
    assert decks.status_code == 200
    assert decks.json()[0]["through_line"] == "Horizon decides the winner."

    missing = await client.get("/reports/nope")
    assert missing.status_code == 404
    assert missing.json()["error"]["code"] == "not_found"


async def test_patch_stores_user_edit_as_new_version(client, sessionmaker, project_with_outputs):
    report_id = project_with_outputs["report_id"]
    response = await client.patch(
        f"/reports/{report_id}", json={"content_markdown": "# Edited by the user\n"}
    )
    assert response.status_code == 200
    body = response.json()
    assert body["version"] == 2
    assert body["id"] != report_id
    assert body["self_check_result"] == {"user_edited": True, "edited_from_version": 1}

    async with sessionmaker() as session:
        original = await session.get(Report, report_id)
        assert "Transformers lead" in original.content_markdown  # v1 untouched
        audited = (
            (
                await session.execute(
                    select(AuditLogEntry).where(
                        AuditLogEntry.action_type == AuditActionType.report_revised.value
                    )
                )
            )
            .scalars()
            .one()
        )
        assert "edited by user" in audited.description


async def test_rewrite_regenerates_for_new_audience_with_provenance(
    client, sessionmaker, event_bus, project_with_outputs
):
    """Acceptance 6: rewrite regenerates for a new audience through the full
    pipeline (outline → draft → self-check) without losing provenance."""
    fake = FakeLLM(
        {
            ReportOutline: outline_responder,
            ReportSection: section_responder,
            SelfCheckResult: [clean_self_check()],
        }
    )

    async def fake_llm_json(messages, schema, **kwargs):
        return fake.complete_json(messages, schema)

    from app.main import app

    app.dependency_overrides[get_llm_json] = lambda: lambda *a, **kw: fake_llm_json(*a, **kw)

    response = await client.post(
        f"/reports/{project_with_outputs['report_id']}/rewrite",
        json={"audience": "executive"},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["version"] == 2
    assert body["audience"] == "executive"
    assert "## Bottom line" in body["content_markdown"]
    assert body["stopping_criterion"] == "saturation"
    assert body["self_check_result"]["clean"] is True

    async with sessionmaker() as session:
        provenance = (
            (await session.execute(select(Provenance).where(Provenance.ref_id == body["id"])))
            .scalars()
            .all()
        )
        # Executive sections: one sourced claim + one flagged inference each.
        assert len(provenance) == 4
        sourced = [p for p in provenance if p.source_id]
        assert len(sourced) == 2
        assert all(p.source_id == project_with_outputs["source_ids"][0] for p in sourced)
        assert all(p.is_inference for p in provenance if not p.source_id)

    ready = event_bus.of_type("output_ready")
    assert len(ready) == 1
    assert ready[0].payload["output"] == "report"
    assert ready[0].payload["report_id"] == body["id"]


async def test_report_export_docx_and_md(client, project_with_outputs):
    report_id = project_with_outputs["report_id"]
    cited = project_with_outputs["source_ids"][0]

    md = await client.get(f"/reports/{report_id}/export", params={"format": "md"})
    assert md.status_code == 200
    assert md.headers["content-disposition"] == 'attachment; filename="report-v1.md"'
    text = md.text
    assert "Transformers lead on short horizons.[1]" in text
    assert "[1] Ada Lovelace (2023). Paper A. Fake Journal." in text
    assert cited not in text  # raw ids replaced by numbers

    docx_response = await client.get(f"/reports/{report_id}/export", params={"format": "docx"})
    assert docx_response.status_code == 200
    assert docx_response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.wordprocessingml"
    )
    document = Document(io.BytesIO(docx_response.content))  # opens uncorrupted
    texts = [p.text for p in document.paragraphs]
    assert "References" in texts

    unsupported = await client.get(f"/reports/{report_id}/export", params={"format": "pdf"})
    assert unsupported.status_code == 422
    assert unsupported.json()["error"]["code"] == "validation_error"


async def test_presentation_export_pptx_and_md(client, project_with_outputs):
    presentation_id = project_with_outputs["presentation_id"]

    pptx_response = await client.get(
        f"/presentations/{presentation_id}/export", params={"format": "pptx"}
    )
    assert pptx_response.status_code == 200
    assert (
        pptx_response.headers["content-disposition"]
        == 'attachment; filename="presentation-v1.pptx"'
    )
    deck = PptxDocument(io.BytesIO(pptx_response.content))  # opens uncorrupted
    titles = [s.shapes.title.text for s in deck.slides]
    assert titles[0] == "Horizon decides the winner."
    assert "References" in titles

    md = await client.get(f"/presentations/{presentation_id}/export", params={"format": "md"})
    assert md.status_code == 200
    assert "**Through-line:** Horizon decides the winner." in md.text
