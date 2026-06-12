"""Export adapters (phase 5 part C): md/docx/pptx generate, open without
corruption, keep confidence labels, and build references lists that match the
inline citations. Exporters are pure transforms — no DB, no fixtures."""

import io
import uuid

import pytest
from docx import Document
from pptx import Presentation as PptxDocument

from app.adapters.export.base import PresentationContent, ReportContent
from app.adapters.export.docx import DocxReportExporter
from app.adapters.export.markdown import ReportMarkdownExporter, SlidesMarkdownExporter
from app.adapters.export.pptx import PptxPresentationExporter
from app.core.errors import ValidationError
from app.db.models import Presentation, Report, Source
from app.services import exports

SID_A = str(uuid.uuid4())
SID_B = str(uuid.uuid4())


def sources() -> dict[str, Source]:
    return {
        SID_A: Source(
            id=SID_A,
            project_id="p",
            title="Paper A",
            authors=["Ada Lovelace"],
            venue="Fake Journal",
            year=2023,
            doi="10.1000/a",
        ),
        SID_B: Source(
            id=SID_B,
            project_id="p",
            title="Paper B",
            authors=["Grace Hopper"],
            year=2024,
            url="https://example.org/b",
        ),
    }


def report_content() -> ReportContent:
    markdown = (
        "# Transformers vs. RNNs\n"
        "\n"
        "*Audience: expert.* Precise hedging.\n"
        "\n"
        "## Findings\n"
        "\n"
        f"Transformers lead on short horizons.[^src:{SID_A}] "
        "*(confidence: well established)*\n"
        "\n"
        f"Long-horizon evidence is thin.[^src:{SID_B}][^src:{SID_A}] "
        "*(confidence: emerging)*\n"
        "\n"
        "- A **bold** bullet with *italic* nuance\n"
    )
    report = Report(
        id="r1", project_id="p", audience="expert", content_markdown=markdown, version=2
    )
    return ReportContent(report=report, sources=sources())


def presentation_content() -> PresentationContent:
    deck = Presentation(
        id="d1",
        project_id="p",
        through_line="Horizon decides the winner.",
        key_messages=[
            {"message": "Short horizons favor transformers.", "source_ids": [SID_A]},
            {"message": "The long game is unproven.", "source_ids": [SID_B]},
            {"message": "Choose per use case.", "source_ids": []},
        ],
        slides=[
            {
                "headline": "Architecture is the wrong question",
                "key_message_index": 0,
                "evidence": [
                    {
                        "text": "0.91 vs 0.84 on short horizons.",
                        "source_ids": [SID_A],
                        "passage": "accuracy of 0.91",
                        "is_inference": False,
                    }
                ],
                "visual": {
                    "type": "comparison_table",
                    "title": "Accuracy by horizon",
                    "columns": ["Model", "Short"],
                    "rows": [["Transformer", "0.91"], ["RNN", "0.84"]],
                    "points": [],
                },
            },
            {
                "headline": "What to do on Monday",
                "key_message_index": 2,
                "evidence": [
                    {
                        "text": "Match model to horizon.",
                        "source_ids": [],
                        "passage": None,
                        "is_inference": True,
                    }
                ],
                "visual": {
                    "type": "bullet_set",
                    "title": None,
                    "columns": [],
                    "rows": [],
                    "points": ["Benchmark on your own horizon"],
                },
            },
        ],
        speaker_notes=[{"slide": 0, "notes": "Caveat: 30-day cap."}, {"slide": 1, "notes": ""}],
        version=1,
    )
    return PresentationContent(presentation=deck, sources=sources())


def test_report_markdown_numbers_citations_and_builds_references():
    text = ReportMarkdownExporter().render(report_content()).decode()
    # Markers renumbered by first appearance; raw ids gone.
    assert "[^src:" not in text
    assert "Transformers lead on short horizons.[1]" in text
    assert "Long-horizon evidence is thin.[2][1]" in text
    # Confidence labels pass through.
    assert "*(confidence: well established)*" in text
    # References match the numbering and the cited sources.
    assert "## References" in text
    assert "[1] Ada Lovelace (2023). Paper A. Fake Journal. https://doi.org/10.1000/a" in text
    assert "[2] Grace Hopper (2024). Paper B. https://example.org/b" in text


def test_report_markdown_export_is_deterministic():
    exporter = ReportMarkdownExporter()
    assert exporter.render(report_content()) == exporter.render(report_content())


def test_docx_opens_with_matching_references_and_labels():
    data = DocxReportExporter().render(report_content())
    document = Document(io.BytesIO(data))  # raises if corrupt
    texts = [p.text for p in document.paragraphs]

    assert "Transformers vs. RNNs" in texts  # heading
    assert "Transformers lead on short horizons.[1] (confidence: well established)" in texts
    assert "Long-horizon evidence is thin.[2][1] (confidence: emerging)" in texts
    # References section matches the cited sources, in citation order.
    refs_at = texts.index("References")
    assert texts[refs_at + 1].startswith("[1] Ada Lovelace (2023). Paper A.")
    assert texts[refs_at + 2].startswith("[2] Grace Hopper (2024). Paper B.")
    # Inline formatting became runs, not literal asterisks.
    assert not any("**" in t for t in texts)
    # Citation markers render superscript.
    findings = next(p for p in document.paragraphs if p.text.startswith("Transformers lead"))
    assert any(r.text == "[1]" and r.font.superscript for r in findings.runs)


def test_pptx_opens_with_slides_visuals_and_notes():
    data = PptxPresentationExporter().render(presentation_content())
    deck = PptxDocument(io.BytesIO(data))  # raises if corrupt
    slides = list(deck.slides)
    titles = [s.shapes.title.text for s in slides]
    assert titles == [
        "Horizon decides the winner.",  # title slide carries the through-line
        "Key messages",
        "Architecture is the wrong question",
        "What to do on Monday",
        "References",
    ]
    # Evidence bullets carry numbered refs; inference is flagged.
    content_slide = slides[2]
    body = content_slide.placeholders[1].text_frame.text
    assert "0.91 vs 0.84 on short horizons. [1]" in body
    assert "Match model to horizon. (inference)" in slides[3].placeholders[1].text_frame.text
    # The visual spec rendered as a native table.
    tables = [sh.table for sh in content_slide.shapes if sh.has_table]
    assert len(tables) == 1
    assert tables[0].cell(0, 0).text == "Model"
    assert tables[0].cell(1, 0).text == "Transformer"
    # Speaker notes land in the notes pane.
    assert content_slide.notes_slide.notes_text_frame.text == "Caveat: 30-day cap."
    # References slide matches the deck-wide numbering.
    refs = slides[4].placeholders[1].text_frame.text
    assert "[1] Ada Lovelace (2023). Paper A." in refs
    assert "[2] Grace Hopper (2024). Paper B." in refs


def test_slides_markdown_fallback():
    text = SlidesMarkdownExporter().render(presentation_content()).decode()
    assert "**Through-line:** Horizon decides the winner." in text
    assert "- Short horizons favor transformers. [1]" in text
    assert "## Slide 1: Architecture is the wrong question" in text
    assert "| Model | Short |" in text
    assert "> Speaker notes: Caveat: 30-day cap." in text
    assert "- Match model to horizon. *(inference)*" in text
    assert "[2] Grace Hopper (2024). Paper B. https://example.org/b" in text


def test_unsupported_formats_are_rejected():
    with pytest.raises(ValidationError):
        exports.report_exporter("pdf")
    with pytest.raises(ValidationError):
        exports.presentation_exporter("docx")
